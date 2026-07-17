from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def pptx_to_markdown(pptx_path, output_md_path, image_dir='Picture'):
    prs = Presentation(pptx_path)
    output = []

    # Ensure image directory exists
    os.makedirs(image_dir, exist_ok=True)

    for i, slide in enumerate(prs.slides, 1):
        output.append(f'# Slide {i}')
        output.append('')

        # Layout name
        layout_name = slide.slide_layout.name if slide.slide_layout else 'Unknown'
        output.append(f'**Layout:** {layout_name}')
        output.append('')

        # Title
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                output.append(f'## {title_text}')
                output.append('')

        # All shapes
        img_count = 0
        for shape in slide.shapes:
            # Skip title (already handled)
            if slide.shapes.title and shape.shape_id == slide.shapes.title.shape_id:
                continue

            # Handle images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_count += 1
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.content_type.split('/')[-1]
                    if ext == 'jpeg':
                        ext = 'jpg'
                    filename = f'slide{i}_img{img_count}.{ext}'
                    filepath = os.path.join(image_dir, filename)
                    with open(filepath, 'wb') as f:
                        f.write(image_bytes)
                    # Add image reference in markdown
                    alt_text = shape.name if shape.name else f'Image {img_count}'
                    output.append(f'![{alt_text}](Picture/{filename})')
                    output.append('')
                except (ValueError, AttributeError):
                    # Linked image or no embedded image
                    output.append(f'[Image: {shape.name} (linked/external)]')
                    output.append('')
                continue

            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # Placeholder info
                ph_info = ''
                if shape.is_placeholder:
                    idx = shape.placeholder_format.idx
                    # Skip slide number placeholder
                    if idx == 12:
                        continue
                    ph_info = f' (placeholder idx={idx})'

                output.append(f'**Shape{ph_info}:**')
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        level = para.level if para.level else 0
                        indent = '  ' * level
                        output.append(f'{indent}- {para_text}')
                output.append('')

            elif shape.has_table:
                table = shape.table
                output.append('**Table:**')
                output.append('')
                # Header row
                header = '| ' + ' | '.join(
                    cell.text.strip().replace('\n', ' ') for cell in table.rows[0].cells
                ) + ' |'
                separator = '| ' + ' | '.join(
                    '---' for _ in table.rows[0].cells
                ) + ' |'
                output.append(header)
                output.append(separator)
                # Data rows
                for row in table.rows[1:]:
                    row_text = '| ' + ' | '.join(
                        cell.text.strip().replace('\n', ' ') for cell in row.cells
                    ) + ' |'
                    output.append(row_text)
                output.append('')

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                output.append('**Speaker Notes:**')
                output.append(f'{notes_text}')
                output.append('')

        output.append('---')
        output.append('')

    return '\n'.join(output)


if __name__ == '__main__':
    import sys

    # Default paths
    pptx_file = r'D:\SourceCode\PROJECTS\CloudMastery\Slide\RAG-on-AWS.pptx'
    output_file = r'D:\SourceCode\PROJECTS\CloudMastery\Slide\RAG-on-AWS.md'
    image_dir = r'D:\SourceCode\PROJECTS\CloudMastery\Slide\Picture'

    # Allow overriding via command line args
    if len(sys.argv) >= 2:
        pptx_file = sys.argv[1]
    if len(sys.argv) >= 3:
        output_file = sys.argv[2]
    if len(sys.argv) >= 4:
        image_dir = sys.argv[3]

    md_content = pptx_to_markdown(pptx_file, output_file, image_dir)
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(md_content)

    print(f'Done! Output: {output_file}')
    print(f'Images saved to: {image_dir}')
